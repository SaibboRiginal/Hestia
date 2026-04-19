import json
import os
import time
import threading
import logging
import io
import tempfile
from urllib.parse import urlparse, parse_qs

import requests

from agents.universal_agent import UniversalAgent
from core.services.context_builder import ContextBuilder
from core.services.memory_service import MemoryService
from core.services.module_registry import ModuleToolRegistry
from core.services.retrieval_service import RetrievalService
from core.services.router_service import RouterService


logger = logging.getLogger(__name__)

# ── Document archiving constants ──────────────────────────────────────────────
_MAX_DOC_ARCHIVE_BYTES = int(
    os.getenv("DOC_MAX_ARCHIVE_BYTES", str(10 * 1024 * 1024)))  # 10 MB
_MAX_EXTRACTED_TEXT_CHARS = 40_000
_CHUNK_SIZE = 900          # characters per chunk
_CHUNK_OVERLAP = 150       # character overlap between consecutive chunks
_MAX_CHUNKS_PER_DOC = 120  # hard cap to avoid runaway embedding costs
_DOC_SEARCH_THRESHOLD = float(
    os.getenv("DOC_SEARCH_THRESHOLD", "1.2"))  # L2 distance

# ── Known multimodal-capable model prefixes ───────────────────────────────────
# Vision (image + PDF): models with visual understanding
_VISION_MODEL_PREFIXES = (
    "gemma-3", "gemma4", "gemma3", "gemma 3", "gemma 4",
    "gemini", "gpt-4o", "gpt-4-vision", "claude-3",
    "llava", "bakllava", "moondream", "minicpm-v", "phi-3-vision", "phi3v",
    "qwen2-vl", "qwen-vl", "internvl",
)
# Audio-native: models that can process raw audio (very few)
_AUDIO_NATIVE_MODEL_PREFIXES = (
    "gemini-1.5", "gemini-2",
)

# ── Lazy singleton loaders for optional local models ─────────────────────────
_CLIP_LOADED = False
_CLIP_MODEL = None
_CLIP_PROCESSOR = None
_CLIP_DEVICE = None


def _load_clip():
    """Lazily load CLIP model once. Returns (model, processor, device) or (None, None, None)."""
    global _CLIP_LOADED, _CLIP_MODEL, _CLIP_PROCESSOR, _CLIP_DEVICE
    if _CLIP_LOADED:
        return _CLIP_MODEL, _CLIP_PROCESSOR, _CLIP_DEVICE
    _CLIP_LOADED = True
    try:
        import torch
        from transformers import CLIPModel, CLIPProcessor as _CLIPProc
        _CLIP_DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
        _CLIP_MODEL = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(_CLIP_DEVICE)
        _CLIP_PROCESSOR = _CLIPProc.from_pretrained("openai/clip-vit-base-patch32")
        _CLIP_MODEL.eval()
        logger.info("[LOCAL] CLIP model loaded (openai/clip-vit-base-patch32) on %s", _CLIP_DEVICE)
    except Exception as exc:
        logger.info("[LOCAL] CLIP unavailable: %s", exc)
    return _CLIP_MODEL, _CLIP_PROCESSOR, _CLIP_DEVICE


_YOLO_LOADED = False
_YOLO_MODEL = None


def _load_yolo():
    """Lazily load YOLOv8-nano once. Returns model or None."""
    global _YOLO_LOADED, _YOLO_MODEL
    if _YOLO_LOADED:
        return _YOLO_MODEL
    _YOLO_LOADED = True
    try:
        from ultralytics import YOLO as _YOLO
        _YOLO_MODEL = _YOLO("yolov8n.pt")
        logger.info("[LOCAL] YOLOv8-nano loaded")
    except Exception as exc:
        logger.info("[LOCAL] YOLO unavailable: %s", exc)
    return _YOLO_MODEL


_WHISPER_LOADED = False
_WHISPER_MODEL = None


def _load_whisper():
    """Lazily load WhisperX base model once. Returns model or None."""
    global _WHISPER_LOADED, _WHISPER_MODEL
    if _WHISPER_LOADED:
        return _WHISPER_MODEL
    _WHISPER_LOADED = True
    try:
        import whisperx
        _WHISPER_MODEL = whisperx.load_model(
            os.getenv("WHISPER_MODEL", "base"),
            device="cpu",
            compute_type="int8",
        )
        logger.info("[LOCAL] WhisperX model loaded (%s)", os.getenv("WHISPER_MODEL", "base"))
    except Exception as exc:
        logger.info("[LOCAL] WhisperX unavailable: %s", exc)
    return _WHISPER_MODEL


class OracleEngine:
    def __init__(self):
        self.hub_api_url = os.getenv(
            "HUB_API_URL", "http://hestia_hub:19001/api").rstrip("/")
        self.archive_url = os.getenv(
            "ARCHIVE_API_URL", "http://hestia_archive:19002/api")

        module_tools_urls_env = os.getenv("MODULE_TOOLS_URLS", "")
        module_tools_url_single = os.getenv("MODULE_TOOLS_URL", "")
        configured_urls = [u.strip()
                           for u in module_tools_urls_env.split(",") if u.strip()]
        if module_tools_url_single:
            configured_urls.append(module_tools_url_single.strip())

        self.context_builder = ContextBuilder(
            max_history_messages=int(os.getenv("ORACLE_HISTORY_LIMIT", "6")),
            max_history_chars=int(
                os.getenv("ORACLE_HISTORY_CHAR_LIMIT", "500")),
            max_entities_in_context=int(
                os.getenv("ORACLE_CONTEXT_ENTITIES_LIMIT", "12")),
            max_field_chars=int(
                os.getenv("ORACLE_CONTEXT_FIELD_CHAR_LIMIT", "280")),
        )

        self.models = {
            "router": {"prov": os.getenv("ROUTER_PROVIDER", "gemini"), "mod": os.getenv("ROUTER_MODEL", "gemma-3-12b-it")},
            "scribe": {
                "prov": os.getenv("SCRIBE_PROVIDER", os.getenv("FALLBACK_ANALYST_PROVIDER", os.getenv("ROUTER_PROVIDER", "ollama"))),
                "mod": os.getenv("SCRIBE_MODEL", os.getenv("FALLBACK_ANALYST_MODEL", os.getenv("ROUTER_MODEL", "qwen2.5:7b"))),
            },
            "analyst": {"prov": os.getenv("ANALYST_PROVIDER", "gemini"), "mod": os.getenv("ANALYST_MODEL", "gemma-3-27b-it")},
            "embedder": {"prov": os.getenv("EMBEDDING_PROVIDER", "ollama"), "mod": os.getenv("EMBEDDING_MODEL", "nomic-embed-text")},
            "fallback_router": {"prov": os.getenv("FALLBACK_ROUTER_PROVIDER", "ollama"), "mod": os.getenv("FALLBACK_ROUTER_MODEL", "mistral:7b")},
            "fallback_scribe": {
                "prov": os.getenv("FALLBACK_SCRIBE_PROVIDER", os.getenv("SCRIBE_PROVIDER", os.getenv("FALLBACK_ROUTER_PROVIDER", "ollama"))),
                "mod": os.getenv("FALLBACK_SCRIBE_MODEL", os.getenv("SCRIBE_MODEL", os.getenv("FALLBACK_ROUTER_MODEL", "mistral:7b"))),
            },
            "fallback_analyst": {"prov": os.getenv("FALLBACK_ANALYST_PROVIDER", "ollama"), "mod": os.getenv("FALLBACK_ANALYST_MODEL", "mistral:7b")},
            "fallback_embedder": {"prov": os.getenv("FALLBACK_EMBEDDING_PROVIDER", "gemini"), "mod": os.getenv("FALLBACK_EMBEDDING_MODEL", "models/embedding-001")},
        }

        self._normalize_provider_model_pairs()

        print(
            f"🧠 Oracle Init | Router: {self.models['router']['mod']} | Scribe: {self.models['scribe']['mod']} | Analyst: {self.models['analyst']['mod']} | Embedder: {self.models['embedder']['mod']}"
        )
        logger.info(
            "Oracle initialized with models | router=%s scribe=%s analyst=%s embedder=%s",
            self.models["router"]["mod"],
            self.models["scribe"]["mod"],
            self.models["analyst"]["mod"],
            self.models["embedder"]["mod"],
        )

        self._init_agents()

        self.module_registry = ModuleToolRegistry(
            module_tool_urls=configured_urls,
            ttl_seconds=int(
                os.getenv("MODULE_TOOL_REGISTRY_TTL_SECONDS", "120")),
            hub_api_url=self.hub_api_url,
        )
        self.router_service = RouterService(self.router, self.fallback_router)
        self.retrieval_service = RetrievalService(
            archive_url=self.archive_url,
            hub_api_url=self.hub_api_url,
            module_registry=self.module_registry,
            embedder=self._embed_text,
        )
        self.memory_service = MemoryService(
            archive_url=self.archive_url,
            hub_api_url=self.hub_api_url,
            scribe_agent=self.scribe,
            fallback_scribe_agent=self.fallback_scribe,
            context_builder=self.context_builder,
        )

    def _conversation_style_contract(self) -> str:
        return """
CONVERSATION STYLE CONTRACT (MANDATORY):
- The reply must feel like an ongoing chat, not a ticket closure.
- Never end with generic assistant closure lines in any language (examples: "Fammi sapere...", "Se in futuro...", "If you need anything else...", "Let me know if...").
- End directly on useful content (fact, answer, suggestion, or next concrete step), without ritual outro.
- Keep tone personal, natural, and context-aware.
""".strip()

    def _normalize_provider_model_pairs(self):
        def normalize(model_key: str, fallback_for_text: str = "gemini-2.5-flash"):
            model_cfg = self.models.get(model_key, {})
            provider = str(model_cfg.get("prov", "")).strip().lower()
            model_name = str(model_cfg.get("mod", "")).strip()

            if provider != "gemini":
                return

            lower_model = model_name.lower()
            if lower_model.startswith("gemma") or ":" in lower_model:
                replacement = "models/embedding-001" if "embed" in model_key else fallback_for_text
                logger.warning(
                    "Invalid Gemini model configured for %s: '%s'. Auto-switching to '%s'.",
                    model_key,
                    model_name,
                    replacement,
                )
                self.models[model_key]["mod"] = replacement

        normalize("router")
        normalize("scribe")
        normalize("analyst")
        normalize("fallback_router")
        normalize("fallback_scribe")
        normalize("fallback_analyst")
        normalize("embedder")
        normalize("fallback_embedder")

    def _init_agents(self):
        router_prompt = """You are a Universal Data Router.
Analyze the user's message and conversation context.

Output ONLY valid JSON with:
1. \"domains\": array of best matching domains. For general chat use [\"general\"].
2. \"filters\": exact-match hints dictionary (or {}).
3. \"filters_gt\": numerical \"greater than\" constraints dictionary (or {}).
4. \"filters_lt\": numerical \"less than\" constraints dictionary (or {}).
5. \"sort_by\": sort field or null.
6. \"sort_order\": \"asc\" or \"desc\".
"""

        scribe_prompt = """You are Hestia's Memory Manager.
Infer enduring user facts/preferences from natural language context.

Do NOT rely on explicit trigger words. Infer intent semantically.
If user asks to forget/reset memory (explicitly or implicitly), output DEPRECATE actions for all active preference IDs.
Output ONLY valid JSON array or NONE.
"""

        analyst_prompt = f"""Sei Hestia, assistente IA universale.

REGOLE CORE:
1. Rispondi nella lingua dell'utente.
2. Usa CONTEXT_DATA_RECORDS solo se pertinente.
3. Applica sempre USER_PREFERENCES.
4. Se i record sono molti, sintetizza e mostra solo i migliori risultati.
5. Puoi attivare notifiche proattive: quando l'utente chiede avvisi/notifiche automatiche, conferma che Hestia può salvarle come sottoscrizioni e inviare alert via Hermes (non dire che non puoi farlo).

FORMATTAZIONE:
- Solo Markdown.
- Usa elenchi puntati e grassetto per chiarezza.
- Se un record contiene URL/link, includilo quando utile.
- LINK POLICY: Per OGNI link/URL, usa SEMPRE il titolo o una descrizione significativa dell'elemento come testo del link.
  ESEMPI CORRETTI: [Appartamento 3 camere via Roma](URL), [Villa con giardino](URL), [Casa luminosa centro](URL)
  MAI USARE: "Apri annuncio", "Clicca qui", "Link", "URL", "Vedi qui" o altri testi generici.
- Non mostrare URL lunghi in chiaro.

STILE FINALE:
{self._conversation_style_contract()}
"""

        self.router = UniversalAgent(
            role_prompt=router_prompt,
            provider=self.models["router"]["prov"],
            model_name=self.models["router"]["mod"],
        )
        self.fallback_router = UniversalAgent(
            role_prompt=router_prompt,
            provider=self.models["fallback_router"]["prov"],
            model_name=self.models["fallback_router"]["mod"],
        )

        self.scribe = UniversalAgent(
            role_prompt=scribe_prompt,
            provider=self.models["scribe"]["prov"],
            model_name=self.models["scribe"]["mod"],
        )
        self.fallback_scribe = UniversalAgent(
            role_prompt=scribe_prompt,
            provider=self.models["fallback_scribe"]["prov"],
            model_name=self.models["fallback_scribe"]["mod"],
        )

        self.analyst = UniversalAgent(
            role_prompt=analyst_prompt,
            provider=self.models["analyst"]["prov"],
            model_name=self.models["analyst"]["mod"],
        )
        self.fallback_analyst = UniversalAgent(
            role_prompt=analyst_prompt,
            provider=self.models["fallback_analyst"]["prov"],
            model_name=self.models["fallback_analyst"]["mod"],
        )

        self.embedder_agent = UniversalAgent(
            role_prompt="",
            provider=self.models["embedder"]["prov"],
            model_name=self.models["embedder"]["mod"],
        )
        self.fallback_embedder_agent = UniversalAgent(
            role_prompt="",
            provider=self.models["fallback_embedder"]["prov"],
            model_name=self.models["fallback_embedder"]["mod"],
        )

    def _emit_status(self, msg: str) -> str:
        return json.dumps({"type": "status", "content": msg}) + "\n"

    def _emit_final(self, reply: str, domain: str = "none") -> str:
        return json.dumps({"type": "final", "reply": reply, "domain": domain}) + "\n"

    def _emit_signal(self, event: str, message: str, data: dict | None = None) -> str:
        return json.dumps(
            {
                "type": "signal",
                "event": event,
                "content": message,
                "data": data or {},
            },
            ensure_ascii=False,
        ) + "\n"

    def _api_get(self, endpoint: str, default_val=None):
        try:
            parsed = urlparse(endpoint)
            endpoint_path = parsed.path if parsed.path.startswith(
                "/") else f"/{parsed.path}"
            query = {k: v[0] if len(v) == 1 else v for k,
                     v in parse_qs(parsed.query).items()}
            normalized = f"api{endpoint_path}"
            response = requests.post(
                f"{self.hub_api_url}/route/archive/{normalized}",
                json={
                    "method": "GET",
                    "query": query,
                    "headers": {},
                    "body": None,
                    "timeout_seconds": 6,
                },
                timeout=7,
            )
            if response.status_code != 200:
                return default_val if default_val is not None else []
            routed = response.json() or {}
            if int(routed.get("status_code", 500)) < 400:
                return routed.get("payload")
            return default_val if default_val is not None else []
        except Exception:
            return default_val if default_val is not None else []

    def _api_post(self, endpoint: str, body: dict, timeout: int = 6):
        normalized = f"api{endpoint if endpoint.startswith('/') else '/' + endpoint}"
        response = requests.post(
            f"{self.hub_api_url}/route/archive/{normalized}",
            json={
                "method": "POST",
                "query": {},
                "headers": {},
                "body": body,
                "timeout_seconds": timeout,
            },
            timeout=timeout + 1,
        )
        response.raise_for_status()
        return response.json() or {}

    def delete_chat_history(self, session_id: str):
        normalized = f"api/chat/history/{session_id}"
        response = requests.post(
            f"{self.hub_api_url}/route/archive/{normalized}",
            json={
                "method": "DELETE",
                "query": {},
                "headers": {},
                "body": None,
                "timeout_seconds": 6,
            },
            timeout=7,
        )
        response.raise_for_status()
        routed = response.json() or {}
        if int(routed.get("status_code", 500)) >= 400:
            raise RuntimeError(routed.get("payload", "delete failed"))
        return routed.get("payload")

    def _embed_text(self, text: str) -> list[float]:
        try:
            vector = self.embedder_agent.embed(text)
            if vector:
                return vector
        except Exception:
            pass

        try:
            vector = self.fallback_embedder_agent.embed(text)
            if vector:
                return vector
        except Exception:
            pass

        return []

    def format_payload(self, command: str, payload: object, response_prompt: str | None = None, client_instructions: str | None = None) -> str:
        payload_text = json.dumps(payload, ensure_ascii=False, indent=2)

        # Detect if this is a proactive alert
        is_alert = str(command or "").startswith("alert:")

        if is_alert:
            formatting_prompt = (
                "Sei Hestia e stai PROATTIVAMENTE informando l'utente. "
                "Scrivi come se TU stessi iniziando una conversazione per condividere qualcosa di rilevante. "
                "Sii naturale, entusiasta ma preciso. "
                "Usa HTML per formattazione (grassetto <b>, link <a href>). "
                "Per i link, usa SEMPRE il titolo/descrizione dell'elemento come testo del link, MAI testi generici. "
                "Non inventare dati. NON usare saluti introduttivi come 'Ciao' o 'Ecco'. "
                f"COMMAND: {command}\n"
                f"SERVICE_PAYLOAD:\n{payload_text}\n"
            )
        else:
            formatting_prompt = (
                "Sei Hestia. Trasforma il payload strutturato in una risposta chiara e utile per l'utente finale. "
                "Mantieni tono naturale, sintetico e orientato all'azione. "
                "Non inventare dati e non includere JSON grezzo se non richiesto. "
                "NON usare saluti, introduzioni o frasi di chiusura rituali. "
                "Rispondi direttamente con i dettagli utili e basta. "
                f"COMMAND: {command}\n"
                f"SERVICE_PAYLOAD:\n{payload_text}\n"
            )

        formatting_prompt += "\n" + self._conversation_style_contract() + "\n"

        if response_prompt and str(response_prompt).strip():
            formatting_prompt += f"\nSERVICE_RESPONSE_PROMPT:\n{str(response_prompt).strip()}\n"

        if client_instructions and str(client_instructions).strip():
            formatting_prompt += f"\nCLIENT_INSTRUCTIONS:\n{str(client_instructions).strip()}\n"

        try:
            return self.analyst.ask(formatting_prompt).strip()
        except Exception:
            return self.fallback_analyst.ask(formatting_prompt).strip()

    def _classify_and_route(
        self,
        user_message: str,
        history_text: str,
        available_domains: list[str],
        schemas: dict | None = None,
    ) -> tuple[str, str | None, float, list[str], dict, dict, dict, str | None, str]:
        domain_candidates = [
            str(domain).strip().lower()
            for domain in (available_domains or [])
            if str(domain).strip().lower() and str(domain).strip().lower() != "general"
        ]
        domain_list_text = ", ".join(
            domain_candidates) if domain_candidates else "none"
        schema_text = json.dumps(
            schemas or {}, ensure_ascii=False, indent=2) if schemas else "{}"

        mode_prompt = f"""You classify and route user intent for a chat orchestrator.

Return ONLY valid JSON with:
1) "mode": "quick_chat" or "domain_query"
2) "domain": one domain from AVAILABLE_DOMAINS or null
3) "confidence": float 0..1
4) "domains": array of routed domains (or ["general"]) for domain_query
5) "filters": exact-match filters object
6) "filters_gt": numeric greater-than filters object
7) "filters_lt": numeric less-than filters object
8) "sort_by": field name or null
9) "sort_order": "asc" or "desc"

Rules:
- Use "quick_chat" for normal conversation, social chat, generic Q&A, short personal exchanges, or messages that do not need structured retrieval.
- Use "domain_query" only when the user clearly asks for domain records, filters, listings, alerts/subscriptions, or data-driven operations.
- Set "domain" only if it is explicit/high-confidence from AVAILABLE_DOMAINS; otherwise null.
- For domain_query, populate domains/filters/sort fields precisely.

AVAILABLE_DOMAINS: {domain_list_text}

CONTEXT DATA STRUCTURES:
{schema_text}

CONTEXT:
{history_text}

USER_MESSAGE: {user_message}
"""

        default_mode = "domain_query"
        default_domain = None
        default_confidence = 0.0
        default_domains = ["general"]
        default_filters: dict = {}
        default_filters_gt: dict = {}
        default_filters_lt: dict = {}
        default_sort_by = None
        default_sort_order = "desc"

        try:
            raw = self.router.ask(mode_prompt).strip()
        except Exception:
            try:
                raw = self.fallback_router.ask(mode_prompt).strip()
            except Exception:
                return (
                    default_mode,
                    default_domain,
                    default_confidence,
                    default_domains,
                    default_filters,
                    default_filters_gt,
                    default_filters_lt,
                    default_sort_by,
                    default_sort_order,
                )

        try:
            start_idx, end_idx = raw.find("{"), raw.rfind("}")
            if start_idx == -1 or end_idx == -1:
                return (
                    default_mode,
                    default_domain,
                    default_confidence,
                    default_domains,
                    default_filters,
                    default_filters_gt,
                    default_filters_lt,
                    default_sort_by,
                    default_sort_order,
                )
            data = json.loads(raw[start_idx: end_idx + 1])

            mode = str(data.get("mode", default_mode)).strip().lower()
            if mode not in {"quick_chat", "domain_query"}:
                mode = default_mode

            domain = data.get("domain")
            normalized_domain = str(domain).strip(
            ).lower() if domain is not None else None
            if normalized_domain and normalized_domain not in domain_candidates:
                normalized_domain = None

            confidence = float(data.get("confidence", 0.0) or 0.0)
            if confidence < 0:
                confidence = 0.0
            if confidence > 1:
                confidence = 1.0

            selected_domains = [str(d).lower() for d in (
                data.get("domains") or []) if str(d).strip()]
            if normalized_domain and normalized_domain not in selected_domains:
                selected_domains.insert(0, normalized_domain)
            valid_domains = [
                d for d in selected_domains if d in available_domains or d == "general"
            ] or ["general"]

            active_filters = data.get("filters") if isinstance(
                data.get("filters"), dict) else {}
            filters_gt = data.get("filters_gt") if isinstance(
                data.get("filters_gt"), dict) else {}
            filters_lt = data.get("filters_lt") if isinstance(
                data.get("filters_lt"), dict) else {}
            sort_by = data.get("sort_by")
            sort_order = "asc" if str(data.get(
                "sort_order", "desc")).lower() == "asc" else "desc"

            return mode, normalized_domain, confidence, valid_domains, active_filters, filters_gt, filters_lt, sort_by, sort_order
        except Exception:
            return (
                default_mode,
                default_domain,
                default_confidence,
                default_domains,
                default_filters,
                default_filters_gt,
                default_filters_lt,
                default_sort_by,
                default_sort_order,
            )

    def _generate_quick_chat_answer(
        self,
        user_message: str,
        history_text: str,
        client_instructions: str | None = None,
        extra_context: str | None = None,
    ) -> str:
        quick_prompt = f"""Sei Hestia, assistente IA conversazionale.

CONTESTO CONVERSAZIONE:
{history_text}
"""
        if extra_context and extra_context.strip():
            quick_prompt += f"\nCONTESTO AGGIUNTIVO:\n{extra_context.strip()}\n"

        quick_prompt += f"""
MESSAGGIO UTENTE: {user_message}

Rispondi in modo naturale, breve (max 3-5 righe), utile e umano.
Se non serve recuperare dati strutturati, resta in conversazione diretta.
"""
        quick_prompt += "\n" + self._conversation_style_contract() + "\n"
        if client_instructions and str(client_instructions).strip():
            quick_prompt += "\n\nSTILE:\n" + str(client_instructions).strip()

        try:
            return self.analyst.ask(quick_prompt)
        except Exception:
            return self.fallback_analyst.ask(quick_prompt)

    def compile_notification_shortcut(self, user_message: str, session_id: str, notify_target: str | None = None) -> dict:
        signals = self.memory_service.extract_and_save_preferences(
            user_message=user_message,
            session_id=session_id,
            notify_target=notify_target,
            force_notification_compiler=True,
        )

        notification_events = {
            "subscription.added",
            "subscription.changed",
            "subscription.removed",
        }
        matched = [
            signal for signal in (signals or [])
            if str(signal.get("event", "")).strip().lower() in notification_events
        ]

        if matched:
            return {
                "ok": True,
                "message": "✅ Notifica elaborata con il comando rapido.",
                "signals": signals,
            }

        return {
            "ok": False,
            "message": "⚠️ Nessuna notifica creata. Specifica meglio dominio, evento o filtri.",
            "signals": signals or [],
        }

    def chat(self, user_message: str, session_id: str, notify_target: str | None = None, force_notification_compiler: bool = False, client_instructions: str | None = None):
        request_started = time.perf_counter()
        logger.info("Chat request started | session_id=%s message_len=%s",
                    session_id, len(user_message or ""))

        # Standard path for data/domain queries
        yield self._emit_status("📂 Recupero cronologia e routing...")

        step_start = time.perf_counter()
        history_data = self._api_get(
            f"/chat/history/{session_id}?limit={self.context_builder.max_history_messages}")
        history_text = self.context_builder.compact_history(history_data)
        logger.info("History loaded | session_id=%s messages=%s in %sms", session_id, len(
            history_data or []), int((time.perf_counter() - step_start) * 1000))

        step_start = time.perf_counter()
        available_domains = self._api_get("/domains") or ["general"]
        logger.info("Domains loaded | count=%s in %sms", len(available_domains), int(
            (time.perf_counter() - step_start) * 1000))

        step_start = time.perf_counter()
        schemas = self._api_get("/schemas") or {}
        logger.info("Metadata loaded | available_domains=%s schema_domains=%s in %sms",
                    available_domains, len(schemas or {}), int((time.perf_counter() - step_start) * 1000))

        classification_start = time.perf_counter()
        mode, explicit_domain, mode_confidence, valid_domains, active_filters, filters_gt, filters_lt, sort_by, sort_order = self._classify_and_route(
            user_message=user_message,
            history_text=history_text,
            available_domains=available_domains,
            schemas=schemas,
        )
        logger.info(
            "Mode/routing classification | mode=%s domain=%s confidence=%.2f in %sms",
            mode,
            explicit_domain,
            mode_confidence,
            int((time.perf_counter() - classification_start) * 1000),
        )

        if mode == "quick_chat" and mode_confidence >= 0.55:
            logger.info(
                "Quick conversation path selected | session_id=%s", session_id)
            yield self._emit_status("💬 Conversazione rapida...")

            # If the user's message is about stored documents, inject the
            # document catalogue so the analyst can answer correctly even in
            # the fast path (no domain_query retrieval happens here).
            doc_context_for_quick = ""
            if self._message_is_about_docs(user_message):
                doc_context_for_quick = self._list_user_docs_brief(
                    chat_id=notify_target, session_id=session_id)
                if doc_context_for_quick:
                    logger.info(
                        "[DOC] Injected brief doc catalogue into quick_chat context")

            analysis_start = time.perf_counter()
            final_answer = self._generate_quick_chat_answer(
                user_message=user_message,
                history_text=history_text,
                client_instructions=client_instructions,
                extra_context=doc_context_for_quick or None,
            )
            logger.info("Quick conversational response generated in %sms", int(
                (time.perf_counter() - analysis_start) * 1000))

            yield self._emit_status("✍️ Consegna...")
            save_start = time.perf_counter()
            try:
                self._api_post(
                    "/chat/history",
                    {"session_id": session_id,
                        "role": "user", "content": user_message},
                )
                self._api_post(
                    "/chat/history",
                    {"session_id": session_id,
                        "role": "assistant", "content": final_answer},
                )
                logger.info("History persisted in %sms", int(
                    (time.perf_counter() - save_start) * 1000))
            except Exception as error:
                logger.warning("Failed persisting chat history: %s", error)

            logger.info("Quick chat request completed | session_id=%s total=%sms",
                        session_id, int((time.perf_counter() - request_started) * 1000))
            yield self._emit_final(final_answer, "general")
            return

        if explicit_domain and explicit_domain not in valid_domains:
            valid_domains = [explicit_domain] + [
                domain for domain in valid_domains if domain != explicit_domain]

        logger.info(
            "Routing complete | domains=%s filters=%s filters_gt=%s filters_lt=%s sort_by=%s sort_order=%s",
            valid_domains,
            active_filters,
            filters_gt,
            filters_lt,
            sort_by,
            sort_order,
        )

        yield self._emit_status(f"🧠 Analisi domini: {', '.join(valid_domains)}...")
        yield self._emit_status("🧾 Recupero preferenze attive...")

        pref_step_start = time.perf_counter()
        all_prefs = []
        seen_pref_ids = set()
        for domain in valid_domains:
            for pref in self._api_get(f"/memory/active?domain={domain}"):
                pref_id = pref.get("id")
                if pref_id and pref_id not in seen_pref_ids:
                    all_prefs.append(pref)
                    seen_pref_ids.add(pref_id)
        logger.info("Preferences loaded | count=%s in %sms", len(
            all_prefs), int((time.perf_counter() - pref_step_start) * 1000))

        preference_facts = [str(p.get("fact", "")).strip()
                            for p in all_prefs if p.get("fact")]

        yield self._emit_status("🔎 Recupero entità dai moduli/Archive...")
        retrieval_start = time.perf_counter()
        all_entities = self.retrieval_service.retrieve_entities(
            user_message=user_message,
            session_id=session_id,
            valid_domains=valid_domains,
            preference_facts=preference_facts,
            active_filters=active_filters,
            filters_gt=filters_gt,
            filters_lt=filters_lt,
            sort_by=sort_by,
            sort_order=sort_order,
        )
        logger.info("Entities retrieval complete | count=%s in %sms", len(
            all_entities), int((time.perf_counter() - retrieval_start) * 1000))

        yield self._emit_status("🧱 Compattazione contesto...")
        formatted_context = self.context_builder.compact_entities_for_prompt(
            all_entities)
        logger.info("Context compacted | has_entities=%s", bool(all_entities))

        # ── Semantic document chunk injection ────────────────────────────────
        # Search stored document chunks semantically relevant to this query.
        # Runs after entity retrieval; uses the same embedding infrastructure.
        doc_chunks = self._search_relevant_docs(
            user_message=user_message,
            chat_id=notify_target,
            session_id=session_id,
        )
        if doc_chunks:
            doc_section = self._format_doc_chunks_for_prompt(doc_chunks)
            formatted_context = (
                (formatted_context + "\n\n" + doc_section).strip()
                if formatted_context
                else doc_section
            )
            logger.info(
                "[DOC] Injected %s chunk(s) from archived docs into context", len(doc_chunks))
        elif self._message_is_about_docs(user_message):
            # No chunks matched but user is asking about their documents —
            # inject the catalogue so the analyst can list them.
            brief = self._list_user_docs_brief(
                chat_id=notify_target, session_id=session_id)
            if brief:
                formatted_context = (
                    (formatted_context + "\n\n" + brief).strip()
                    if formatted_context
                    else brief
                )
                logger.info(
                    "[DOC] Injected brief doc catalogue (no chunk match) into domain_query context")

        analysis_prompt = self.context_builder.build_analysis_prompt(
            preference_facts=preference_facts,
            valid_domains=valid_domains,
            active_filters=active_filters,
            filters_gt=filters_gt,
            filters_lt=filters_lt,
            sort_by=sort_by,
            sort_order=sort_order,
            formatted_context=formatted_context,
            history_text=history_text,
            user_message=user_message,
        )
        analysis_prompt += "\n\n" + self._conversation_style_contract()

        if client_instructions and str(client_instructions).strip():
            analysis_prompt += (
                "\n\nCLIENT_INSTRUCTIONS:\n"
                + str(client_instructions).strip()
            )

        yield self._emit_status("🧠 Sintesi finale in corso...")
        analysis_start = time.perf_counter()
        try:
            final_answer = self.analyst.ask(analysis_prompt)
            logger.info("Analyst response generated with primary model in %sms", int(
                (time.perf_counter() - analysis_start) * 1000))
        except Exception as primary_error:
            logger.warning(
                "Primary analyst failed, switching to fallback: %s", primary_error)
            fallback_start = time.perf_counter()
            try:
                final_answer = self.fallback_analyst.ask(analysis_prompt)
                logger.info("Analyst response generated with fallback model in %sms", int(
                    (time.perf_counter() - fallback_start) * 1000))
            except Exception as fallback_error:
                logger.error("Fallback analyst failed: %s", fallback_error)
                final_answer = (
                    "⚠️ In questo momento i modelli sono temporaneamente non disponibili. "
                    "Riprova tra poco."
                )

        yield self._emit_status("✍️ Consegna...")

        save_start = time.perf_counter()
        try:
            self._api_post(
                "/chat/history",
                {"session_id": session_id,
                    "role": "user", "content": user_message},
            )
            self._api_post(
                "/chat/history",
                {"session_id": session_id,
                    "role": "assistant", "content": final_answer},
            )
            logger.info("History persisted in %sms", int(
                (time.perf_counter() - save_start) * 1000))
        except Exception as error:
            logger.warning("Failed persisting chat history: %s", error)
            pass

        yield self._emit_status("🔔 Aggiornamento preferenze e notifiche...")
        memory_start = time.perf_counter()
        try:
            signals = self.memory_service.extract_and_save_preferences(
                user_message,
                session_id,
                notify_target=notify_target,
                force_notification_compiler=force_notification_compiler,
            )
            logger.info(
                "Preference/subscription sync completed | signals=%s in %sms",
                len(signals or []),
                int((time.perf_counter() - memory_start) * 1000),
            )
            for signal in signals or []:
                yield self._emit_signal(
                    event=str(signal.get("event", "info")),
                    message=str(signal.get(
                        "message", "Aggiornamento eseguito.")),
                    data=signal.get("data") or {},
                )
        except Exception as error:
            logger.warning("Preference/subscription sync failed: %s", error)

        logger.info("Chat request completed | session_id=%s total=%sms",
                    session_id, int((time.perf_counter() - request_started) * 1000))
        yield self._emit_final(final_answer, valid_domains[0])

    def extract_and_save_preferences(self, user_message: str, session_id: str):
        self.memory_service.extract_and_save_preferences(
            user_message, session_id)

    # ── Document archiving helpers ────────────────────────────────────────────

    def _analyst_supports_vision(self) -> bool:
        """Return True if the current analyst model is known to support image/PDF vision input."""
        model = (self.analyst.model_name or "").lower().replace("_", "-")
        return any(model.startswith(pfx.lower().replace("_", "-")) for pfx in _VISION_MODEL_PREFIXES)

    def _analyst_supports_audio(self) -> bool:
        """Return True if the current analyst model supports raw audio input natively."""
        model = (self.analyst.model_name or "").lower()
        return any(model.startswith(pfx.lower()) for pfx in _AUDIO_NATIVE_MODEL_PREFIXES)

    # ── Static text-extraction helpers (no LLM needed) ───────────────────────

    @staticmethod
    def _extract_plain_text(file_bytes: bytes) -> str:
        """Decode plain text files (txt, csv, log, md, rst…) with encoding fallback."""
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                return file_bytes.decode(enc)
            except (UnicodeDecodeError, ValueError):
                continue
        return file_bytes.decode("utf-8", errors="replace")

    @staticmethod
    def _extract_json(file_bytes: bytes) -> str:
        """Pretty-print JSON for readable embedding."""
        try:
            data = json.loads(file_bytes.decode("utf-8", errors="replace"))
            return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception:
            return file_bytes.decode("utf-8", errors="replace")

    @staticmethod
    def _extract_html(file_bytes: bytes) -> str:
        """Strip HTML tags and return visible text."""
        try:
            from bs4 import BeautifulSoup
            return BeautifulSoup(
                file_bytes.decode("utf-8", errors="replace"), "html.parser"
            ).get_text(separator=" ", strip=True)
        except ImportError:
            # Naive strip if BeautifulSoup not installed
            import re
            raw = file_bytes.decode("utf-8", errors="replace")
            return re.sub(r"<[^>]+>", " ", raw)

    @staticmethod
    def _extract_docx(file_bytes: bytes) -> str:
        """Extract text from .docx using python-docx."""
        try:
            import docx as _docx
            doc = _docx.Document(io.BytesIO(file_bytes))
            parts: list[str] = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            # Also extract table cells
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        parts.append(row_text)
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("[DOC] docx extraction failed: %s", exc)
            return ""

    @staticmethod
    def _extract_odf(file_bytes: bytes) -> str:
        """Extract text from .odt / .ods / .odp using odfpy."""
        try:
            from odf.opendocument import load as _odf_load
            from odf import text as _odf_text, teletype as _odf_tt
            doc = _odf_load(io.BytesIO(file_bytes))
            texts = doc.getElementsByType(_odf_text.P)
            return "\n".join(_odf_tt.extractText(t) for t in texts if _odf_tt.extractText(t).strip())
        except Exception as exc:
            logger.warning("[DOC] odf extraction failed: %s", exc)
            return ""

    @staticmethod
    def _extract_xlsx(file_bytes: bytes) -> str:
        """Extract text from .xlsx using openpyxl."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            parts: list[str] = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("[DOC] xlsx extraction failed: %s", exc)
            return ""

    @staticmethod
    def _extract_xls(file_bytes: bytes) -> str:
        """Extract text from legacy .xls using xlrd."""
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=file_bytes)
            parts: list[str] = []
            for sheet in wb.sheets():
                parts.append(f"[Sheet: {sheet.name}]")
                for r in range(sheet.nrows):
                    row_text = " | ".join(str(sheet.cell_value(r, c)) for c in range(sheet.ncols))
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("[DOC] xls extraction failed: %s", exc)
            return ""

    @staticmethod
    def _extract_pptx(file_bytes: bytes) -> str:
        """Extract text from .pptx using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {slide_num}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("[DOC] pptx extraction failed: %s", exc)
            return ""

    @staticmethod
    def _extract_pdf_text(file_bytes: bytes) -> str:
        """Extract text from PDF using pypdf (no LLM)."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            parts = []
            for page in reader.pages:
                t = page.extract_text()
                if t and t.strip():
                    parts.append(t)
            return "\n".join(parts)
        except Exception as exc:
            logger.warning("[DOC] pypdf text extraction failed: %s", exc)
            return ""

    # ── Local vision analysis: CLIP + YOLO ───────────────────────────────────

    @staticmethod
    def _analyze_image_local(file_bytes: bytes) -> dict:
        """Run CLIP zero-shot classification + YOLO object detection on an image.

        Returns a dict with:
          - ``description``: natural language description built from detections
          - ``tags``: list of detected labels / top CLIP categories
          - ``clip_available``: bool
          - ``yolo_available``: bool
        """
        from PIL import Image as _PILImage

        try:
            pil_img = _PILImage.open(io.BytesIO(file_bytes)).convert("RGB")
        except Exception as exc:
            logger.warning("[DOC] Cannot open image for local analysis: %s", exc)
            return {"description": "", "tags": [], "clip_available": False, "yolo_available": False}

        tags: list[str] = []
        yolo_available = False
        clip_available = False
        yolo_lines: list[str] = []
        clip_lines: list[str] = []

        # YOLO object detection
        yolo = _load_yolo()
        if yolo is not None:
            yolo_available = True
            try:
                results = yolo(pil_img, verbose=False)
                detected: dict[str, int] = {}
                for result in results:
                    for box in result.boxes:
                        cls_name = result.names[int(box.cls[0])].replace("_", " ")
                        detected[cls_name] = detected.get(cls_name, 0) + 1
                if detected:
                    for obj, count in sorted(detected.items(), key=lambda x: -x[1]):
                        tags.append(obj)
                        yolo_lines.append(f"  - {count}× {obj}")
            except Exception as exc:
                logger.warning("[DOC] YOLO inference failed: %s", exc)

        # CLIP zero-shot classification against a broad scene/object taxonomy
        clip_model, clip_proc, clip_device = _load_clip()
        if clip_model is not None and clip_proc is not None:
            clip_available = True
            try:
                import torch
                _CLIP_LABELS = [
                    "a photo of a person", "a photo of an animal", "a photo of a vehicle",
                    "a photo of a building or architecture", "a photo of nature or landscape",
                    "a photo of food or drink", "a photo of text or document", "a photo of art or painting",
                    "a photo of electronics or technology", "a photo of furniture or interior",
                    "a photo of sport or fitness activity", "a photo of medical or scientific content",
                    "a photo of a map or diagram", "a photo of a chart or graph",
                    "a screenshot of an interface or application",
                ]
                inputs = clip_proc(text=_CLIP_LABELS, images=pil_img, return_tensors="pt", padding=True)
                inputs = {k: v.to(clip_device) for k, v in inputs.items()}
                with torch.no_grad():
                    outputs = clip_model(**inputs)
                    probs = outputs.logits_per_image.softmax(dim=1)[0]
                top_idxs = probs.argsort(descending=True)[:3]
                for idx in top_idxs:
                    if float(probs[idx]) > 0.05:
                        label = _CLIP_LABELS[idx].replace("a photo of ", "").replace("a screenshot of ", "")
                        clip_lines.append(f"  - {label} ({probs[idx]:.0%})")
                        # Add short keyword to tags if not already there
                        short = label.split(" or ")[0].split(" and ")[0].strip()
                        if short not in tags:
                            tags.append(short)
            except Exception as exc:
                logger.warning("[DOC] CLIP inference failed: %s", exc)

        desc_parts: list[str] = []
        if yolo_lines:
            desc_parts.append("Detected objects:\n" + "\n".join(yolo_lines))
        if clip_lines:
            desc_parts.append("Visual scene categories:\n" + "\n".join(clip_lines))

        return {
            "description": "\n".join(desc_parts),
            "tags": tags[:10],
            "clip_available": clip_available,
            "yolo_available": yolo_available,
        }

    # ── Audio transcription: WhisperX ────────────────────────────────────────

    @staticmethod
    def _transcribe_audio(file_bytes: bytes, mime_type: str) -> str:
        """Transcribe audio (or audio track of video) using WhisperX.

        Returns transcribed text or empty string if unavailable / failed.
        """
        # Determine file suffix from mime type
        _MIME_TO_EXT = {
            "audio/mpeg": ".mp3", "audio/mp3": ".mp3",
            "audio/wav": ".wav", "audio/x-wav": ".wav",
            "audio/ogg": ".ogg", "audio/vorbis": ".ogg",
            "audio/flac": ".flac",
            "audio/aac": ".aac", "audio/x-aac": ".aac",
            "audio/m4a": ".m4a", "audio/mp4": ".m4a",
            "video/mp4": ".mp4", "video/mpeg": ".mpeg",
            "video/webm": ".webm", "video/ogg": ".ogv",
            "video/quicktime": ".mov", "video/x-msvideo": ".avi",
        }
        ext = _MIME_TO_EXT.get(mime_type, ".audio")
        whisper = _load_whisper()
        if whisper is None:
            return ""
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as f:
                f.write(file_bytes)
                tmp_path = f.name
            result = whisper.transcribe(tmp_path, batch_size=8)
            segments = result.get("segments") or []
            return " ".join(seg.get("text", "").strip() for seg in segments if seg.get("text", "").strip())
        except Exception as exc:
            logger.warning("[DOC] WhisperX transcription failed: %s", exc)
            return ""
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    @staticmethod
    def _chunk_text(text: str) -> list[str]:
        """Split text into overlapping chunks for fine-grained RAG retrieval."""
        text = text.strip()
        if not text:
            return []
        chunks: list[str] = []
        start = 0
        while start < len(text):
            end = min(start + _CHUNK_SIZE, len(text))
            chunks.append(text[start:end])
            if end == len(text):
                break
            start = end - _CHUNK_OVERLAP
        return chunks[:_MAX_CHUNKS_PER_DOC]

    def _extract_and_archive_document(
        self,
        file_bytes: bytes,
        mime_type: str,
        final_answer: str,
        document_id: str,
        session_id: str,
        chat_id: str | None,
        filename: str | None,
    ) -> None:
        """Background thread: extract full text/metadata, chunk, embed, save to Archive.

        Routing logic (in order of precedence):
          audio / video  → WhisperX transcription → (optionally) LLM for metadata
          image          → YOLO+CLIP local analysis → LLM vision (if capable) for richer metadata
          PDF            → LLM multimodal (if capable) else pypdf text extraction
          office docs    → python-docx / odfpy / openpyxl / python-pptx text extraction
          text / code    → direct decode
          unknown        → fallback to analysis answer already produced by Oracle
        """
        import hashlib as _hashlib

        # SHA-256 for dedup and future blob-storage association
        file_hash = _hashlib.sha256(file_bytes).hexdigest()

        # Capability flags for the current analyst model
        model_has_vision = self._analyst_supports_vision()
        model_has_audio = self._analyst_supports_audio()

        # Fetch available Hestia domains so the LLM can pick the best one
        try:
            available_domains: list[str] = self._api_get("/domains") or []
        except Exception:
            available_domains = []
        domains_hint = ", ".join(available_domains) if available_domains else "documents, general"

        extracted_text = ""
        title: str | None = None
        summary: str | None = None
        domain: str = "documents"
        tags: list[str] = []

        # LLM extraction prompt (for vision-capable paths)
        def _llm_extraction_prompt(text_field_hint: str = "complete extracted text") -> str:
            return (
                "Analyze this content thoroughly.\n"
                "Output ONLY valid JSON with exactly these keys:\n"
                '{"title": "<concise title, max 120 chars>", '
                '"summary": "<2-3 sentence summary, max 400 chars>", '
                f'"text": "<{text_field_hint}>", '
                f'"domain": "<best match from: {domains_hint}>", '
                '"tags": ["<kw1>", "<kw2>", "<kw3>"]}\n'
                "tags: 3-8 short keywords. No commentary. No markdown fences. Only JSON."
            )

        def _parse_llm_json(raw: str) -> dict:
            s, e = raw.find("{"), raw.rfind("}") + 1
            if s >= 0 and e > s:
                return json.loads(raw[s:e])
            return {}

        def _apply_llm_data(data: dict):
            nonlocal title, summary, extracted_text, domain, tags
            title = data.get("title") or title
            summary = data.get("summary") or summary
            extracted_text = data.get("text") or extracted_text
            if data.get("domain"):
                domain = str(data["domain"]).strip().lower()
            if data.get("tags"):
                tags = [str(t) for t in data["tags"] if str(t).strip()]

        def _llm_ask_with_attachment(prompt: str) -> str:
            """Try primary analyst then fallback for multimodal attachment calls."""
            try:
                return self.analyst.ask_with_attachment(
                    file_bytes=file_bytes, mime_type=mime_type, user_message=prompt)
            except Exception as exc1:
                logger.warning("[DOC] Primary analyst attachment call failed: %s", exc1)
                try:
                    return self.fallback_analyst.ask_with_attachment(
                        file_bytes=file_bytes, mime_type=mime_type, user_message=prompt)
                except Exception as exc2:
                    logger.warning("[DOC] Fallback analyst attachment call failed: %s", exc2)
                    return ""

        # ──────────────────────────────────────────────────────────────────────
        # BRANCH 1 — AUDIO & VIDEO → WhisperX transcription
        # ──────────────────────────────────────────────────────────────────────
        is_audio = mime_type.startswith("audio/")
        is_video = mime_type.startswith("video/")
        if is_audio or is_video:
            logger.info("[DOC] Audio/video path | model_audio=%s", model_has_audio)
            transcribed = ""
            # Try WhisperX first (works offline, highly accurate)
            transcribed = self._transcribe_audio(file_bytes, mime_type)

            if not transcribed and model_has_audio:
                # Native audio model (e.g. Gemini 1.5): let LLM transcribe
                prompt = _llm_extraction_prompt("complete verbatim transcription of the audio/video")
                raw = _llm_ask_with_attachment(prompt)
                if raw:
                    try:
                        _apply_llm_data(_parse_llm_json(raw))
                    except Exception:
                        transcribed = raw[:_MAX_EXTRACTED_TEXT_CHARS]

            if transcribed and not extracted_text:
                extracted_text = transcribed

            # If we have transcribed text but no metadata yet, ask a text-only LLM to enrich
            if extracted_text and not title:
                meta_prompt = (
                    f"Given this audio transcript, produce only valid JSON:\n"
                    f'{{"title": "...", "summary": "...", "domain": "<from: {domains_hint}>", "tags": [...]}}\n'
                    f"Transcript:\n{extracted_text[:2000]}"
                )
                try:
                    raw_meta = self.analyst.ask(meta_prompt)
                    _apply_llm_data(_parse_llm_json(raw_meta))
                except Exception:
                    pass

        # ──────────────────────────────────────────────────────────────────────
        # BRANCH 2 — IMAGES → CLIP+YOLO local, then LLM vision if capable
        # ──────────────────────────────────────────────────────────────────────
        elif mime_type.startswith("image/"):
            logger.info("[DOC] Image path | model_vision=%s", model_has_vision)

            # Always run local analysis first (fast, offline, enriches tags)
            try:
                local = self._analyze_image_local(file_bytes)
                if local["tags"]:
                    tags = local["tags"]
                if local["description"]:
                    extracted_text = local["description"]
                logger.info(
                    "[DOC] Local image analysis | CLIP=%s YOLO=%s tags=%s",
                    local["clip_available"], local["yolo_available"], tags,
                )
            except Exception as exc:
                logger.warning("[DOC] Local image analysis error: %s", exc)

            # If the analyst supports vision, get a richer LLM description
            if model_has_vision:
                prompt = _llm_extraction_prompt("complete visual description of the image")
                raw = _llm_ask_with_attachment(prompt)
                if raw:
                    try:
                        data = _parse_llm_json(raw)
                        # Merge LLM tags with local YOLO/CLIP tags
                        llm_tags = [str(t) for t in (data.get("tags") or []) if str(t).strip()]
                        merged_tags = list(dict.fromkeys(llm_tags + tags))[:10]
                        data["tags"] = merged_tags
                        _apply_llm_data(data)
                    except Exception as exc:
                        logger.warning("[DOC] Image LLM parse failed: %s", exc)
                        if not extracted_text:
                            extracted_text = raw[:_MAX_EXTRACTED_TEXT_CHARS]

            # If no LLM or LLM failed, synthesise text from local detections
            if not extracted_text and tags:
                extracted_text = "Image content: " + ", ".join(tags)

        # ──────────────────────────────────────────────────────────────────────
        # BRANCH 3 — PDF → LLM multimodal (if capable) else pypdf text
        # ──────────────────────────────────────────────────────────────────────
        elif mime_type == "application/pdf":
            logger.info("[DOC] PDF path | model_vision=%s", model_has_vision)
            if model_has_vision:
                prompt = _llm_extraction_prompt("complete verbatim text content of the PDF")
                raw = _llm_ask_with_attachment(prompt)
                if raw:
                    try:
                        _apply_llm_data(_parse_llm_json(raw))
                    except Exception:
                        extracted_text = extracted_text or raw[:_MAX_EXTRACTED_TEXT_CHARS]

            # Always supplement with pypdf if LLM text came out empty
            if not extracted_text:
                logger.info("[DOC] PDF fallback: using pypdf text extraction")
                extracted_text = self._extract_pdf_text(file_bytes)

        # ──────────────────────────────────────────────────────────────────────
        # BRANCH 4 — OFFICE DOCUMENTS (Word / LibreOffice / Excel / PowerPoint)
        # ──────────────────────────────────────────────────────────────────────
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ):
            logger.info("[DOC] DOCX/DOC path")
            extracted_text = self._extract_docx(file_bytes)

        elif mime_type in (
            "application/vnd.oasis.opendocument.text",
            "application/vnd.oasis.opendocument.spreadsheet",
            "application/vnd.oasis.opendocument.presentation",
        ):
            logger.info("[DOC] ODF path")
            extracted_text = self._extract_odf(file_bytes)

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ):
            logger.info("[DOC] XLSX path")
            extracted_text = self._extract_xlsx(file_bytes)

        elif mime_type in ("application/vnd.ms-excel",):
            logger.info("[DOC] XLS path")
            extracted_text = self._extract_xls(file_bytes)

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ):
            logger.info("[DOC] PPTX path")
            extracted_text = self._extract_pptx(file_bytes)

        # ──────────────────────────────────────────────────────────────────────
        # BRANCH 5 — TEXT-BASED (plain, CSV, JSON, HTML, Markdown, code…)
        # ──────────────────────────────────────────────────────────────────────
        elif mime_type.startswith("text/") or mime_type in (
            "application/json",
            "application/x-yaml", "application/yaml",
            "application/xml", "text/xml",
        ):
            logger.info("[DOC] Text path | mime=%s", mime_type)
            if "json" in mime_type:
                extracted_text = self._extract_json(file_bytes)
            elif "html" in mime_type:
                extracted_text = self._extract_html(file_bytes)
            else:
                extracted_text = self._extract_plain_text(file_bytes)

        # ──────────────────────────────────────────────────────────────────────
        # BRANCH 6 — UNKNOWN / BINARY → fallback to Oracle's analysis answer
        # ──────────────────────────────────────────────────────────────────────
        else:
            logger.info("[DOC] Unknown mime type '%s', using final_answer as text", mime_type)

        # For office/text docs with no metadata yet, ask a text-only LLM to enrich
        if extracted_text and not title and not is_audio and not is_video and not mime_type.startswith("image/"):
            meta_prompt = (
                "Given the following document content, produce only valid JSON:\n"
                f'{{"title": "...", "summary": "...", "domain": "<best match from: {domains_hint}>", "tags": [...]}}\n'
                "tags: 3-8 short keywords.\n"
                f"Content (first 2000 chars):\n{extracted_text[:2000]}"
            )
            try:
                raw_meta = self.analyst.ask(meta_prompt)
                _apply_llm_data(_parse_llm_json(raw_meta))
            except Exception as exc:
                logger.debug("[DOC] Metadata enrichment LLM call failed: %s", exc)

        # ── Fallback defaults ─────────────────────────────────────────────────
        if not extracted_text:
            extracted_text = final_answer
        if not summary:
            summary = (extracted_text[:400] or final_answer[:400]).strip() or None

        # Normalise domain against known Hestia domains
        if available_domains and domain not in available_domains:
            domain = "documents"

        # ── Chunking + embedding ──────────────────────────────────────────────
        extracted_text = extracted_text[:_MAX_EXTRACTED_TEXT_CHARS]
        chunks = self._chunk_text(extracted_text)
        context_prefix = f"Document: {title or filename or 'Attachment'}\n{summary or ''}\n\n"

        embedded_chunks = []
        for i, chunk_text in enumerate(chunks):
            emb = self._embed_text(context_prefix + chunk_text)
            embedded_chunks.append({
                "chunk_index": i,
                "chunk_text": chunk_text,
                "embedding": emb if emb else None,
            })

        doc_emb_text = f"{title or ''}\n{summary or ''}\n{extracted_text[:400]}".strip()
        doc_embedding = self._embed_text(doc_emb_text) if doc_emb_text else []

        body = {
            "document_id": document_id,
            "session_id": session_id,
            "chat_id": chat_id,
            "filename": filename,
            "mime_type": mime_type,
            "file_size_bytes": len(file_bytes),
            "file_hash": file_hash,
            "title": title,
            "summary": summary,
            "extracted_text": extracted_text,
            "embedding": doc_embedding if doc_embedding else None,
            "is_permanent": False,
            "domain": domain,
            "tags": json.dumps(tags) if tags else None,
            "chunks": embedded_chunks,
        }
        try:
            self._api_post("/documents", body, timeout=60)
            logger.info(
                "[DOC] Archived | id=%s chunks=%s title=%r domain=%s tags=%s model_vision=%s",
                document_id, len(embedded_chunks), title, domain, tags, model_has_vision,
            )
        except Exception as exc:
            logger.warning("[DOC] Archive save failed: %s", exc)

    def _search_relevant_docs(
        self,
        user_message: str,
        chat_id: str | None,
        session_id: str | None,
        limit: int = 4,
    ) -> list[dict]:
        """Semantic search over stored document chunks relevant to the current query."""
        query_vector = self._embed_text(user_message)
        if not query_vector:
            return []
        try:
            body = {
                "query_vector": query_vector,
                "chat_id": chat_id,
                "session_id": session_id,
                "limit": limit,
                "threshold": _DOC_SEARCH_THRESHOLD,
            }
            routed = self._api_post("/documents/search", body, timeout=8)
            payload = routed.get("payload") if isinstance(
                routed, dict) else routed
            if isinstance(payload, list):
                return payload
        except Exception as exc:
            logger.debug("[DOC] Chunk search failed (non-fatal): %s", exc)
        return []

    @staticmethod
    @staticmethod
    def _format_doc_chunks_for_prompt(chunks: list[dict]) -> str:
        """Format retrieved document chunks into a prompt section."""
        if not chunks:
            return ""
        lines = ["📎 RELEVANT ARCHIVED DOCUMENT PASSAGES:"]
        seen: dict[str, str] = {}
        for chunk in chunks:
            doc_id = chunk.get("document_id", "?")
            title = chunk.get("title") or chunk.get(
                "filename") or "Untitled document"
            if doc_id not in seen:
                seen[doc_id] = title
                lines.append(f"\n[{title}]")
            lines.append(f"  …{chunk.get('chunk_text', '')}…")
        return "\n".join(lines)

    def _list_user_docs_brief(self, chat_id: str | None, session_id: str | None, limit: int = 10) -> str:
        """Return a compact catalogue of the user's stored documents for injection into the prompt.

        Used when the user asks something like "what documents have you saved?" so that
        the analyst has an explicit list even without a semantic chunk match.
        """
        try:
            query: dict = {"limit": str(limit)}
            if chat_id:
                query["chat_id"] = str(chat_id)
            elif session_id:
                query["session_id"] = str(session_id)

            resp = requests.post(
                f"{self.hub_api_url}/route/archive/api/documents",
                json={"method": "GET", "query": query,
                      "headers": {}, "body": {}, "timeout_seconds": 6},
                timeout=7,
            )
            if resp.status_code != 200:
                return ""
            docs: list[dict] = resp.json().get("payload") or []
            if not docs:
                return ""

            lines = [f"📎 USER'S ARCHIVED DOCUMENTS ({len(docs)} stored):"]
            for doc in docs:
                title = doc.get("title") or doc.get("filename") or "Untitled"
                domain = doc.get("domain", "documents")
                perm = "📌 permanent" if doc.get(
                    "is_permanent") else "temporary"
                accessed = doc.get("access_count", 0)
                tags_raw = doc.get("tags")
                tag_str = ""
                if tags_raw:
                    try:
                        import json as _json
                        tag_list = _json.loads(tags_raw) if isinstance(
                            tags_raw, str) else tags_raw
                        if isinstance(tag_list, list):
                            tag_str = f" [{', '.join(str(t) for t in tag_list[:4])}]"
                    except Exception:
                        pass
                lines.append(
                    f"  • {title}{tag_str} — domain:{domain}, {perm}, accessed {accessed}×")
            return "\n".join(lines)
        except Exception as exc:
            logger.debug("[DOC] Brief doc list failed (non-fatal): %s", exc)
            return ""

    _DOC_AWARENESS_KEYWORDS = frozenset([
        "document", "documents", "documento", "documenti",
        "file", "files",
        "pdf", "attachment", "allegato",
        "saved", "salvato", "salvati", "archiviato", "archiviati",
        "uploaded", "caricato", "caricati",
        "hai", "have", "know about", "remember",
    ])

    def _message_is_about_docs(self, message: str) -> bool:
        """Heuristic: does this message likely ask about stored documents?"""
        lower = message.lower()
        return any(kw in lower for kw in self._DOC_AWARENESS_KEYWORDS)

    def analyze_document(
        self,
        file_bytes: bytes,
        mime_type: str,
        user_message: str,
        session_id: str,
        notify_target: str | None = None,
        client_instructions: str | None = None,
        filename: str | None = None,
    ):
        """Reason over an attached file and stream an NDJSON reply.

        Supported: images, PDFs, audio/video, office docs, text files.
        Falls back gracefully based on model capabilities.
        """
        is_audio = mime_type.startswith("audio/")
        is_video = mime_type.startswith("video/")
        model_has_vision = self._analyst_supports_vision()
        model_has_audio = self._analyst_supports_audio()
        # A doc the current model cannot see natively
        needs_local_extraction = (is_audio or is_video) and not model_has_audio

        if is_audio or is_video:
            yield self._emit_status("🎙️ Trascrizione audio in corso...")
        else:
            yield self._emit_status("📄 Analisi documento in corso...")

        full_prompt = (
            f"{self._conversation_style_contract()}\n\n"
            f"The user has attached a file and is asking the following:\n"
            f"{user_message}"
        )
        if client_instructions and str(client_instructions).strip():
            full_prompt += f"\n\nCLIENT_INSTRUCTIONS:\n{str(client_instructions).strip()}"

        final_answer = ""

        if needs_local_extraction:
            # For audio/video with a non-audio-capable model: transcribe locally
            # then answer the user's question as a text query.
            yield self._emit_status("🔊 Trascrizione con WhisperX...")
            transcript = self._transcribe_audio(file_bytes, mime_type)
            if transcript:
                text_prompt = (
                    f"{self._conversation_style_contract()}\n\n"
                    f"The user sent an audio/video file. Here is its transcript:\n"
                    f"---\n{transcript[:6000]}\n---\n\n"
                    f"User's question: {user_message}"
                )
                if client_instructions and str(client_instructions).strip():
                    text_prompt += f"\n\nCLIENT_INSTRUCTIONS:\n{str(client_instructions).strip()}"
                try:
                    yield self._emit_status("🤖 Analisi trascrizione...")
                    final_answer = self.analyst.ask(text_prompt)
                except Exception:
                    try:
                        final_answer = self.fallback_analyst.ask(text_prompt)
                    except Exception:
                        final_answer = f"📝 Trascrizione:\n\n{transcript[:2000]}"
            else:
                final_answer = "⚠️ Non riesco a trascrivere l'audio. Prova a inviare un file in un formato supportato (mp3, wav, ogg, m4a)."
        else:
            # For all other files (images, PDF, docs) or audio-capable model
            try:
                yield self._emit_status("🤖 Elaborazione con LLM...")
                final_answer = self.analyst.ask_with_attachment(
                    file_bytes=file_bytes,
                    mime_type=mime_type,
                    user_message=full_prompt,
                )
            except Exception as primary_exc:
                logger.warning("Primary analyst document analysis failed: %s", primary_exc)
                try:
                    final_answer = self.fallback_analyst.ask_with_attachment(
                        file_bytes=file_bytes,
                        mime_type=mime_type,
                        user_message=full_prompt,
                    )
                except Exception as fallback_exc:
                    logger.error("Fallback analyst also failed: %s", fallback_exc)
                    # Last resort for non-vision model receiving a text-based doc:
                    # extract text locally and ask as text
                    local_text = ""
                    if mime_type == "application/pdf":
                        local_text = self._extract_pdf_text(file_bytes)
                    elif "word" in mime_type or "docx" in mime_type:
                        local_text = self._extract_docx(file_bytes)
                    elif mime_type.startswith("text/") or "json" in mime_type:
                        local_text = self._extract_plain_text(file_bytes)
                    if local_text:
                        fallback_text_prompt = (
                            f"{self._conversation_style_contract()}\n\n"
                            f"Document content:\n---\n{local_text[:6000]}\n---\n\n"
                            f"User question: {user_message}"
                        )
                        try:
                            final_answer = self.analyst.ask(fallback_text_prompt)
                        except Exception:
                            final_answer = "⚠️ Non riesco ad analizzare il documento in questo momento."
                    else:
                        final_answer = "⚠️ Non riesco ad analizzare il documento in questo momento."

        # Persist this turn in chat history
        try:
            self._api_post("/chat/history", {
                "session_id": session_id,
                "role": "user",
                "content": f"[File: {filename or mime_type}] {user_message}",
            })
            self._api_post("/chat/history", {
                "session_id": session_id,
                "role": "assistant",
                "content": final_answer,
            })
        except Exception as hist_exc:
            logger.warning("Failed to persist document analysis history: %s", hist_exc)

        # ── Archive document for RAG ──────────────────────────────────────────
        import uuid as _uuid
        document_id = _uuid.uuid4().hex
        should_archive = len(file_bytes) <= _MAX_DOC_ARCHIVE_BYTES

        if should_archive:
            _file_bytes = file_bytes  # capture for closure
            threading.Thread(
                target=self._extract_and_archive_document,
                args=(
                    _file_bytes, mime_type, final_answer,
                    document_id, session_id, notify_target, filename,
                ),
                daemon=True,
            ).start()
            yield self._emit_signal(
                event="document_saved",
                message="📎 Documento salvato nell'archivio.",
                data={
                    "document_id": document_id,
                    "filename": filename,
                    "mime_type": mime_type,
                    "file_size_bytes": len(file_bytes),
                },
            )
        else:
            logger.info(
                "[DOC] File too large to archive (%s bytes > %s limit), analysis only.",
                len(file_bytes), _MAX_DOC_ARCHIVE_BYTES,
            )

        yield self._emit_final(final_answer, "document")
