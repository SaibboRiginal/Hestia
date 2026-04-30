"""Pure file-type text extractors — no LLM, no I/O side effects.

Single responsibility: given raw bytes and a MIME type, return plain text.

Open/Closed: register new extractors in _EXTRACTOR_REGISTRY without touching
existing extractor functions or the dispatch logic.
"""
import io
import json
import logging
from typing import Callable

logger = logging.getLogger(f"hestia_oracle.{__name__}")

# ── Individual extractors ─────────────────────────────────────────────────────


def _plain_text(file_bytes: bytes) -> str:
    """Decode plain text with encoding fallback (utf-8, latin-1, cp1252)."""
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            return file_bytes.decode(enc)
        except (UnicodeDecodeError, ValueError):
            continue
    return file_bytes.decode("utf-8", errors="replace")


def _json_text(file_bytes: bytes) -> str:
    """Pretty-print JSON for readable embedding."""
    try:
        data = json.loads(file_bytes.decode("utf-8", errors="replace"))
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        return file_bytes.decode("utf-8", errors="replace")


def _html_text(file_bytes: bytes) -> str:
    """Strip HTML tags; falls back to a naive regex if BeautifulSoup is absent."""
    try:
        from bs4 import BeautifulSoup
        return BeautifulSoup(
            file_bytes.decode("utf-8", errors="replace"), "html.parser"
        ).get_text(separator=" ", strip=True)
    except ImportError:
        import re
        raw = file_bytes.decode("utf-8", errors="replace")
        return re.sub(r"<[^>]+>", " ", raw)


def _pdf_text(file_bytes: bytes) -> str:
    """Extract text from PDF using pypdf (no LLM)."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        parts = [p.extract_text() for p in reader.pages if p.extract_text()]
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("event=extract_pypdf_failed [EXTRACT] pypdf failed: %s", exc)
        return ""


def _docx_text(file_bytes: bytes) -> str:
    """Extract text from .docx (paragraphs + table cells)."""
    try:
        import docx as _docx
        doc = _docx.Document(io.BytesIO(file_bytes))
        parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip()
                                      for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("event=extract_docx_failed [EXTRACT] docx failed: %s", exc)
        return ""


def _odf_text(file_bytes: bytes) -> str:
    """Extract text from .odt / .ods / .odp using odfpy."""
    try:
        from odf.opendocument import load as _load
        from odf import text as _t, teletype as _tt
        doc = _load(io.BytesIO(file_bytes))
        return "\n".join(
            _tt.extractText(p)
            for p in doc.getElementsByType(_t.P)
            if _tt.extractText(p).strip()
        )
    except Exception as exc:
        logger.warning("event=extract_odf_failed [EXTRACT] odf failed: %s", exc)
        return ""


def _xlsx_text(file_bytes: bytes) -> str:
    """Extract text from .xlsx using openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(
            file_bytes), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("event=extract_xlsx_failed [EXTRACT] xlsx failed: %s", exc)
        return ""


def _xls_text(file_bytes: bytes) -> str:
    """Extract text from legacy .xls using xlrd."""
    try:
        import xlrd
        wb = xlrd.open_workbook(file_contents=file_bytes)
        parts: list[str] = []
        for sheet in wb.sheets():
            parts.append(f"[Sheet: {sheet.name}]")
            for r in range(sheet.nrows):
                row_text = " | ".join(str(sheet.cell_value(r, c))
                                      for c in range(sheet.ncols))
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("event=extract_xls_failed [EXTRACT] xls failed: %s", exc)
        return ""


def _pptx_text(file_bytes: bytes) -> str:
    """Extract text from .pptx slide by slide."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("event=extract_pptx_failed [EXTRACT] pptx failed: %s", exc)
        return ""


# ── MIME → extractor registry (Open/Closed: add here, nothing else changes) ──

_EXTRACTOR_REGISTRY: dict[str, Callable[[bytes], str]] = {
    # PDF
    "application/pdf": _pdf_text,
    # Word
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _docx_text,
    "application/msword": _docx_text,
    # ODF
    "application/vnd.oasis.opendocument.text": _odf_text,
    "application/vnd.oasis.opendocument.spreadsheet": _odf_text,
    "application/vnd.oasis.opendocument.presentation": _odf_text,
    # Excel
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _xlsx_text,
    "application/vnd.ms-excel": _xls_text,
    # PowerPoint
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _pptx_text,
    # Structured text
    "application/json": _json_text,
    "application/x-yaml": _plain_text,
    "application/yaml": _plain_text,
    "application/xml": _plain_text,
    "text/xml": _plain_text,
}


def extract_text(file_bytes: bytes, mime_type: str) -> str:
    """Dispatch *file_bytes* to the appropriate extractor for *mime_type*.

    Returns extracted plain text, or an empty string if no extractor matched
    and the content is binary/unknown.
    """
    # Exact MIME match
    extractor = _EXTRACTOR_REGISTRY.get(mime_type)
    if extractor:
        return extractor(file_bytes)

    # Prefix-based fallbacks for text/* types
    if "html" in mime_type:
        return _html_text(file_bytes)
    if mime_type.startswith("text/"):
        return _plain_text(file_bytes)

    # No match — return empty (caller decides the fallback)
    return ""
